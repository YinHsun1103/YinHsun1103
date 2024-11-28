
import streamlit as st  # 匯入 Streamlit 庫，用於創建網頁應用
import random  # 匯入 random 庫，提供隨機數生成功能
import pandas as pd  # 匯入 Pandas 庫，用於資料處理與分析
import numpy as np  # 匯入 NumPy 庫，用於數值計算
from PIL import Image  # 從 PIL 庫匯入 Image 模組，用於處理圖像
import datetime  # 匯入 datetime 庫，用於時間和日期的操作
import openai  # 匯入 OpenAI 庫，用於與 OpenAI 的 API 進行互動，例如 GPT 模型的使用
import os      # 匯入 os 庫，用於操作系統層級的功能，例如環境變數管理、文件路徑操作等

import requests


# 設置頁面配置
st.set_page_config(page_title="ChatGPT-like clone")

#-----------------------------------------------------------------------------------------------------------------------------------------

# 定義頁籤選項
# 此處建立了一個包含三個頁籤名稱的清單，每個頁籤代表不同的網頁內容
tabs = ["Steamlit練習", "HomeWork1", "HomeWork2", "HomeWork3", "HomeWork4"]

# 初始化 session state
# 檢查 "selected_tab" 是否已在 session_state 中，若無則進行初始化
if "selected_tab" not in st.session_state:
    st.session_state.selected_tab = tabs[0]  # 若尚未設定，預設選中第一個頁籤

# 使用 sidebar 的 radio 選項讓使用者選擇頁面
# 在側邊欄中創建一個 radio 按鈕，讓使用者可以選擇不同的頁籤
selected_tab = st.sidebar.radio("選擇頁面", tabs)

# 更新 session_state 的頁籤狀態
# 將選擇的頁籤更新到 session_state 中，確保頁面狀態的一致性
st.session_state.selected_tab = selected_tab

# 根據選定的頁籤顯示相對應的頁面內容
# 檢查當前選中的頁籤，並根據選中的頁籤顯示相應的內容
if st.session_state.selected_tab == "Steamlit練習":
    st.title("Steamlit練習")  # 在頁面中顯示標題為 "Steamlit練習"

    # 設定網頁標題，顯示文字 'Hello, streamlit！我的第一支web應用程式開發!!!'
    st.title('Hello, streamlit！我的第一支web應用程式開發!!!')

    # 使用 st.write 方法顯示一個簡單的字串
    st.write('我是一個字串')

    # 定義一個數字變數 K 並顯示
    K = 9999
    st.write(K)

    # 建立一個 DataFrame，內容為隨機數據，包含 10 列和 20 行，並且每列標頭為 "col_0" 到 "col_19"
    dataframe = pd.DataFrame(np.random.randn(10, 20), columns=[f"col_{i}" for i in range(20)])

    # 將 DataFrame 顯示在頁面上
    st.write(dataframe)

    # 使用 markdown 顯示標題，示範文字格式化的能力
    st.markdown("# 這是一個示例 Streamlit 網頁")
    st.markdown('''
        :red[Streamlit] :orange[can] :green[write] :blue[text] :violet[in]
        :gray[pretty] :rainbow[colors].''')
    # 顯示一行包含表情符號的文本
    st.markdown("Here's a bouquet &mdash; #:tulip::cherry_blossom::rose::hibiscus::sunflower::blossom:")

    # 設定頁面標題，顯示歡迎文字
    st.title("歡迎來到我的網站")
    # 顯示頁面標題，並展示文字的格式化功能
    st.title('_Streamlit_ is :blue[cool] :sunglasses:')
    # 顯示標題
    st.header("這是一個標頭", divider='rainbow')
    # 顯示帶有文字和表情符號的標頭
    st.header('_Streamlit_ is :blue[cool] :sunglasses:')
    # 顯示子標頭，並加上分隔線
    st.subheader("這是一個子標頭", divider='rainbow')
    st.subheader('_Streamlit_ is :blue[cool] :sunglasses:')

    # 使用 caption 顯示圖片的描述文字
    st.caption("這是一張美麗的圖片")
    st.caption('A caption with italics :blue[colors] and emojis :sunglasses:')

    # 顯示 Python 程式碼段落
    st.code("print('Hello, Streamlit!')")
    # 定義一段多行程式碼並顯示，啟用行號顯示
    code = '''def hello():
        print("Hello, Streamlit!")'''
    st.code(code, line_numbers=True)    
    st.code(code, language='python')

    # 顯示純文字內容
    st.text("這是一些純文字內容。")

    # 使用 LaTeX 語法顯示數學公式
    st.latex(r"e^{i\pi} + 1 = 0")
    # 顯示一個更複雜的公式，包含求和和分數表示
    st.latex(r'''
        a + ar + a r^2 + a r^3 + \cdots + a r^{n-1} =
        \sum_{k=0}^{n-1} ar^k =
        a \left(\frac{1-r^{n}}{1-r}\right)
        ''')

    # 插入分隔線
    st.divider()
    # 在分隔線之間插入一段文字
    st.write("慘了，我被夾在分隔線中間！")
    # 再次插入分隔線
    st.divider()

    # 建立另一個 DataFrame，包含隨機數據，並顯示它
    df = pd.DataFrame(np.random.randn(10, 10), columns=("col %d" % i for i in range(10)))
    st.dataframe(df)  # 與 st.write(df) 相同
    # 顯示 DataFrame，並將每列的最大值高亮顯示
    st.dataframe(df.style.highlight_max(axis=0))

    # 建立一個包含應用資訊的 DataFrame，包含應用名稱、網址、GitHub 星星數和過去 30 天的瀏覽歷史數據
    df = pd.DataFrame(
        {
            "name": ["Roadmap", "Extras", "Issues"],
            "url": ["https://roadmap.streamlit.app", "https://extras.streamlit.app", "https://issues.streamlit.app"],
            "stars": [random.randint(0, 1000) for _ in range(3)],  # 隨機生成 0 到 1000 的星星數
            "views_history": [[random.randint(0, 5000) for _ in range(30)] for _ in range(3)],  # 每天隨機生成 0 到 5000 的瀏覽數據
        }
    )

    # 在頁面上顯示 DataFrame，並對列進行自訂設定
    st.dataframe(
        df,
        column_config={
            "name": "App name",  # 自訂欄位名稱為 "App name"
            "stars": st.column_config.NumberColumn(
                "Github Stars",  # 顯示 GitHub 星星數的欄位
                help="Number of stars on GitHub",  # 顯示提示訊息
                format="%d ⭐",  # 格式化為整數並加上星星符號
            ),
            "url": st.column_config.LinkColumn("App URL"),  # 將 URL 欄位顯示為可點擊連結
            "views_history": st.column_config.LineChartColumn(
                "Views (past 30 days)", y_min=0, y_max=5000  # 使用折線圖顯示瀏覽數據，並設置 y 軸範圍
            ),
        },
        hide_index=True  # 隱藏索引欄位
    )
    use_container_width = True  # 設置 DataFrame 使用全頁面寬度

    # 顯示一些數據指標 (metrics)
    st.metric(label="溫度", value="30 °C", delta="1.2 °C")  # 顯示溫度及其變動情況
    
    # 使用三欄佈局分別顯示不同的數據指標
    col1, col2, col3 = st.columns(3)
    col1.metric("溫度", "30 °C", "1.2 °C")  # 第一欄顯示溫度
    col2.metric("風力", "9 mph", "-8%")  # 第二欄顯示風速
    col3.metric("濕度", "86%", "4%")  # 第三欄顯示濕度

    # 顯示股票和金價數據，並設定增減顏色的變化
    st.metric(label="金價", value=3580, delta=-250, delta_color="inverse")
    st.metric(label="聯發科", value=1100, delta=80, delta_color="inverse")
    st.metric(label="台積電", value=512, delta=0, delta_color="off")

    # JSON 資料顯示，展示使用字典儲存的資料
    data = {
        '姓名': '王小明',
        '年齡': 30,
        '地址': '台北市',
        '學歷': {
            '學士學位': '資訊科學',
            '碩士學位': '資訊管理',
        },
        '興趣': [
            '運動',
            '讀書',
            '旅遊',
        ],
    }
    st.json(data)  # 將 JSON 資料顯示為格式化的結構

    # 顯示圖片，並添加說明文字
    image = Image.open('S__24829963.jpg')  # 從檔案中讀取圖片
    st.image(image, caption='這是一隻阿拉斯加的照片')  # 顯示圖片，並設置說明文字

    # 建立並顯示面積圖
    chart_data = pd.DataFrame(np.random.randn(20, 3), columns=["a", "b", "c"])
    st.area_chart(chart_data)  # 顯示面積圖

    # 建立並顯示長條圖
    st.bar_chart(chart_data)  # 顯示長條圖

    # 建立並顯示折線圖
    st.line_chart(chart_data)  # 顯示折線圖

    # 建立並顯示散點圖
    st.scatter_chart(chart_data)  # 顯示散點圖

    # 添加變更大小的散點圖，根據 col3 大小變化
    chart_data = pd.DataFrame(np.random.randn(20, 3), columns=["col1", "col2", "col3"])
    chart_data['col4'] = np.random.choice(['A', 'B', 'C'], 20)  # 添加分類列
    st.scatter_chart(chart_data, x='col1', y='col2', color='col4', size='col3')  # 顯示不同大小與顏色的散點圖

    # 顯示地圖，隨機生成經緯度數據，模擬地圖數據
    df = pd.DataFrame(np.random.randn(1000, 2) / [50, 50] + [37.76, -122.4], columns=['lat', 'lon'])
    st.map(df)  # 顯示地圖，將經緯度資料標記在地圖上

    # 定義台北市的各區域邊界資訊
    taipei_boundaries = {
        "區域": ["中正區", "大同區", "中山區", "松山區", "大安區", "萬華區", "信義區", "士林區", "北投區", "內湖區", "南港區", "文山區"],
        "西經度": [121.506623, 121.511034, 121.527879, 121.566173, 121.552743, 121.500144, 121.578663, 121.521708, 121.667708, 121.588943, 121.618678, 121.570432],
        "東經度": [121.529235, 121.525196, 121.570859, 121.583734, 121.577419, 121.518360, 121.592992, 121.540610, 121.529424, 121.639907, 121.632538, 121.611367],
        "南緯度": [25.032883, 25.062731, 25.069360, 25.048741, 25.026515, 25.035154, 25.031934, 25.089020, 25.110381, 25.073524, 25.044769, 24.989247],
        "北緯度": [25.041144, 25.066231, 25.080802, 25.050934, 25.041457, 25.046256, 25.050930, 25.146468, 25.153252, 25.123361, 25.091840, 25.006014]
    }

    df_boundaries = pd.DataFrame(taipei_boundaries)

    # 生成隨機數據，模擬各區域的人口數
    df_boundaries['人口數'] = np.random.randint(100, 1000, len(df_boundaries))

    # 為每個區域創建一個隨機顏色，用於地圖顯示
    df_boundaries['顏色'] = ['#{:02x}{:02x}{:02x}'.format(np.random.randint(0, 256), np.random.randint(0, 256), np.random.randint(0, 256)) for _ in range(len(df_boundaries))]

    # 顯示地圖，標示台北市各區域，並使用人口數決定標記大小，顏色依據隨機顏色
    st.map(df_boundaries,
           latitude='南緯度',
           longitude='西經度',
           size='人口數',
           color='顏色'
    )

    # 按鈕區域設置
    st.button("重設", type="primary")
    if st.button(":100:我是一個按鈕"):
        st.write("現在按下的是「我是一個按鈕」！！！")
    if st.button("這是第二個按鈕"):
        st.write("現在按下的是「第二個按鈕」！！！")
    if st.button(":dart:這是第二個按鈕"):
        st.write("現在按下的是，帶有圖示的「第二個按鈕」！！！")

    # 連結按鈕
    st.link_button("前往奇摩首頁", "https://tw.yahoo.com", type="primary", help="hello my friend")
    st.link_button("前往台中科大首頁", "https://www.nutc.edu.tw/", disabled=True)

    # 下載按鈕 - 下載資料框為 CSV
    data = {
        'Column1': [1, 2, 3, 4, 5],
        'Column2': ['A', 'B', 'C', 'D', 'E']
    }
    my_large_df = pd.DataFrame(data)

    @st.cache_data
    def convert_df(df):
        # 將 DataFrame 轉為 CSV 格式，並使用 cache 減少重新計算次數
        return df.to_csv(index=False).encode('utf-8')

    csv = convert_df(my_large_df)

    st.download_button(
        label="dataframe下載成csv",
        data=csv,
        file_name='large_df.csv',
        mime='text/csv'
    )

    # 單純文字下載，存為 txt 檔
    text_contents = '''這就是單純的將文字下載，並存成txt檔'''
    st.download_button('將文字下載', text_contents)

    # 下載圖片
    with open("S__24829963.jpg", "rb") as file:
        btn = st.download_button(
            label="下載圖片",
            data=file,
            file_name="阿拉斯加.jpg",
            mime="image/png"
        )

    # 電影標題輸入應用
    st.title('電影標題輸入應用')
    st.write('請在下方輸入電影標題，然後按下確認按鈕。')
    movie_title = st.text_input('輸入電影標題', '今夜作學也會笑', max_chars=15)
    if st.button('確認電影標題'):
        st.write('您選擇的電影標題是：', movie_title)

    # 密碼輸入應用
    st.title('密碼輸入應用')
    st.write('請在下方輸入您的密碼，然後按下確認按鈕。')
    password = st.text_input('輸入密碼', type='password')
    if st.button('確認密碼'):
        # 檢查密碼是否正確
        if password == '88888888':  # 將'88888888'替換為您的實際密碼
            st.write('密碼正確！歡迎您進入系統。')
        else:
            st.write('密碼錯誤，請檢查後重新輸入。')

    # 創意文字分析應用
    st.title('創意文字分析應用')
    st.write('在下方的文字區域中，輸入您的奇幻故事，然後按下確認按鈕。')
    story_text = st.text_area('輸入您的奇幻故事', max_chars=100)
    if st.button("確認"):
        if story_text:
            words = story_text.split(" ")
            word_count = len(words)
            st.write(f'您的故事中包含了 {word_count} 段文字！')
        else:
            st.write('請先輸入您的奇幻故事再按下確認按鈕。')

    # 數字輸入
    number = st.number_input('輸入一個數字', value=None, placeholder='請在這裡輸入...', min_value=0, max_value=100)
    st.write('您輸入的是:', number)

    # 生日輸入應用
    st.title('生日輸入應用')
    st.write('請輸入您的生日日期，然後按下確認按鈕。')
    birthday = st.date_input("您的生日日期", datetime.date(1990, 1, 1))
    if st.button('確認生日'):
        st.write('您的生日是：', birthday.strftime('%Y年%m月%d日'))

    # 選擇明年休假期間
    today = datetime.datetime.now()
    next_year = today.year + 1
    jan_1 = datetime.date(next_year, 1, 1)
    dec_31 = datetime.date(next_year, 12, 31)
    d = st.date_input(
        "選擇明年要休假的期間",
        (jan_1, datetime.date(next_year, 1, 7)),
        min_value=jan_1,
        max_value=dec_31,
        format="YYYY-MM-DD"
    )
    d

    # 設定自動發信的時間
    t = st.time_input('設定自動發信的時間', value=None, step=3600)
    st.write('自動發信時間', t)

    # 建立三欄佈局
    col1, col2, col3 = st.columns(3)

    # 顯示阿拉斯加圖片於第一欄
    with col1:
       st.header("阿拉斯加")  # 設定標題
       st.image("S__24829963.jpg")  # 顯示阿拉斯加的圖片

    # 顯示邊境牧羊犬圖片於第二欄
    with col2:
       st.header("邊境牧羊犬")  # 設定標題
       st.image("images.jpg")  # 顯示邊境牧羊犬的圖片

    # 顯示柴犬圖片於第三欄
    with col3:
       st.header("柴犬")  # 設定標題
       st.image("1a0ec01465964e9fa986689864e47f3d_th.jpg.crdownload")  # 顯示柴犬的圖片

    # 建立三個分頁標籤
    tab1, tab2, tab3 = st.tabs(["阿拉斯加", "邊境牧羊犬", "柴犬"])

    # 在第一個分頁中顯示阿拉斯加圖片
    with tab1:
        st.header("阿拉斯加")  # 設定標題
        st.image("S__24829963.jpg")  # 顯示阿拉斯加的圖片

    # 在第二個分頁中顯示邊境牧羊犬圖片
    with tab2:
        st.header("邊境牧羊犬")  # 設定標題
        st.image("images.jpg")  # 顯示邊境牧羊犬的圖片

    # 在第三個分頁中顯示柴犬圖片
    with tab3:
        st.header("柴犬")  # 設定標題
        st.image("1a0ec01465964e9fa986689864e47f3d_th.jpg.crdownload")  # 顯示柴犬的圖片

# ---------------------------------------------------------------------------------------------------------------
    # HomeWork1分頁
elif st.session_state.selected_tab == "HomeWork1":
    st.title("可上傳CSV/WORD/PDF/EXCEL 檔")  # 設定標題
    st.write("請上傳您的文件.")  # 顯示提示文字

    # 文件上傳功能，限制文件格式為 CSV, DOCX, PDF, EXCEL
    file = st.file_uploader("選擇文件", type=['csv', 'docx', 'pdf', 'xlsx'])

    if file is not None:
        file_type = file.name.split('.')[-1].lower()

        if file_type == 'csv':
            # 處理 CSV 文件
            df = pd.read_csv(file)
            st.write("以下是您上傳的 CSV 數據：")
            st.write(df)

            selected_column = st.selectbox("選擇要繪製的列", df.columns)

            st.subheader("折線圖")
            st.line_chart(df[selected_column])

        elif file_type == 'xlsx':
            # 處理 Excel 文件
            df = pd.read_excel(file)
            st.write("以下是您上傳的 EXCEL 數據：")
            st.write(df)

            selected_column = st.selectbox("選擇要繪製的列", df.columns)

            st.subheader("折線圖")
            st.line_chart(df[selected_column])

        elif file_type == 'docx':
            # 處理 Word 文件，直接使用 file 來讀取
            doc = docx.Document(file)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            st.write("以下是您上傳的 Word 文件內容：")
            st.write("\n".join(full_text))

        elif file_type == 'pdf':
            # 處理 PDF 文件
            pdf_reader = PyPDF2.PdfReader(file)
            pdf_text = ""
            for page_num in range(len(pdf_reader.pages)):
                pdf_text += pdf_reader.pages[page_num].extract_text()
            st.write("以下是您上傳的 PDF 文件內容：")
            st.write(pdf_text)
#-----------------------------------------------------------------------------------------------------------------------------------------
elif st.session_state.selected_tab == "HomeWork2":
    st.title("HomeWork2 之 連接Chat GPT")  # 顯示標題

    # 設置 API 金鑰
    openai.api_key = "sk-svcacct-fb_-GzpFTmE6wtv222EkZdGrZrVUnZdTIP-AkvTvtcxO8n7D-tZvHHAL6ChEGT3BlbkFJCwdg-PbyzjyhbVo99UJNUKYTHayGD-I0QpeVibX_K7x6F8UE9Q7j0flr-VmAA"  # 從 Streamlit 秘密配置中取得 OpenAI API 金鑰
    openai.api_base = "https://api.openai.com/v1/chat/completions"  # 設定 OpenAI 的 API 金鑰，使應用可以呼叫 OpenAI API


    # 初始化應用狀態
    if "openai_model" not in st.session_state:
        st.session_state["openai_model"] = "gpt-3.5-turbo"  # 設定預設使用的聊天模型（例如 gpt-3.5-turbo）

    if "messages" not in st.session_state:
        st.session_state.messages = []  # 初始化訊息列表，存放聊天記錄

    # 顯示整個聊天紀錄
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):  # 根據訊息的角色顯示（使用者或助手）
            st.markdown(message["content"])  # 以 markdown 格式顯示訊息內容

    # 聊天輸入框
    if prompt := st.chat_input("What is up?"):  # 聊天輸入框，用戶可輸入訊息
        # 添加用戶輸入的訊息到訊息列表
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):  # 顯示用戶的訊息
            st.markdown(prompt)  # 以 markdown 格式顯示用戶的訊息

        # 呼叫 OpenAI API 以獲取助手回應
        with st.chat_message("assistant"):  # 顯示助手回應的訊息框
            try:
                response = openai.ChatCompletion.create(
                    model=st.session_state["openai_model"],  # 使用已設置的聊天模型
                    messages=st.session_state.messages  # 將整個訊息記錄發送至 OpenAI API
                )
                assistant_reply = response['choices'][0]['message']['content']  # 取得助手的回應內容
                # 添加助手的回應到訊息列表
                st.session_state.messages.append({"role": "assistant", "content": assistant_reply})
                st.markdown(assistant_reply)  # 以 markdown 格式顯示助手的回應

            except Exception as e:
                st.error(f"An error occurred: {e}")  # 若發生錯誤，顯示錯誤訊息

    # 聊天記錄上限提示
    if len(st.session_state.messages) >= 20:  # 設定訊息數量上限（此處為 20）
        st.info(
            """Notice: The maximum message limit for this demo version has been reached. 
            We encourage you to build your own application using Streamlit's tutorial."""
        )  # 當訊息達到上限時，顯示提示訊息


# ---------------------------------------------------------------------------------------------------------------
    # HomeWork3分頁
elif st.session_state.selected_tab == "HomeWork3":
    # 加入自訂樣式的 HTML 和 CSS
    st.markdown("""
        <style>
            .main-title {
                background-color: #4CAF50;
                color: white;
                padding: 15px;
                border-radius: 10px;
                text-align: center;
            }
            div.stButton > button {
                background-color: #FF6347;
                color: white;
                font-size: 16px;
                padding: 10px 20px;
                border-radius: 10px;
            }
            div.stFileUploader {
                border: 2px dashed #4CAF50;
                padding: 15px;
                border-radius: 10px;
                margin: 15px 0;
            }
            .message {
                margin: 10px 0;
                padding: 15px;
                border-radius: 10px;
                box-shadow: 2px 2px 10px rgba(0, 0, 0, 0.1);
            }
            .user-message {
                background-color: #f0f0f0;
            }
            .ai-message {
                background-color: #d4f7d4;
            }
        </style>
    """, unsafe_allow_html=True)

    # 標題區
    st.markdown("""
        <div class="main-title">
            <h1>可上傳檔案讓 AI 分析</h1>
        </div>
    """, unsafe_allow_html=True)
    # 初始化 Streamlit 介面
    st.write("可上傳一個檔案(txt.pdf.docx.csv.pptx)讓GPT為您做分析")
    

    # 使用提供的 ChatGPT Anywhere API Key
    api_key = "sk-UebxWPd44CJnjYPWt3e85ogp7DOqXSj2AuGEUarE9NjUckni"
    api_url = "https://api.chatanywhere.tech/v1/chat/completions"

    # 初始化 session_state 來儲存聊天記錄
    if "messages" not in st.session_state:
        st.session_state["messages"] = [{"role": "system", "content": "你是一個幫助解答問題的助手，請用繁體中文回答。"}]

    # 上傳檔案，支援 txt, pdf, docx, csv, pptx
    uploaded_file = st.file_uploader("上傳檔案", type=["txt", "pdf", "docx", "csv", "pptx"])

    # 使用者輸入問題
    user_input = st.text_area("請輸入您的問題：", placeholder="輸入訊息後按下 '送出'")

    # 定義訊息樣式
    def format_user_message(message):
        return f"""
        <div style='background-color: #d3d3d3; padding: 10px; border-radius: 20px; margin: 10px 0;'>
            <strong>你：</strong> {message}
        </div>
        """

    def format_ai_message(message):
        return f"""
        <div style='background-color: #d4f7d4; padding: 10px; border-radius: 20px; margin: 10px 0;'>
            <strong>AI：</strong> {message}
        </div>
        """

    # 讀取上傳檔案的內容，移除 py 的處理邏輯
    def read_uploaded_file(file):
        if file is not None:
            try:
                # 根據檔案類型處理
                if file.type == "text/plain":  # txt 檔案
                    return file.read().decode("utf-8")
                elif file.type == "application/pdf":
                    import PyPDF2
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in range(len(pdf_reader.pages)):
                        text += pdf_reader.pages[page].extract_text()
                    return text
                elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    import docx
                    doc = docx.Document(file)
                    text = "\n".join([para.text for para in doc.paragraphs])
                    return text
                elif file.type == "text/csv":
                    df = pd.read_csv(file)
                    return df.to_string()  # 轉成可讀取的文字
                elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":  # pptx
                    from pptx import Presentation
                    prs = Presentation(file)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                else:
                    return "無法解析該檔案格式。"
            except Exception as e:
                return f"讀取檔案時發生錯誤：{e}"
        return None

    # 分析上傳檔案的內容
    file_content = read_uploaded_file(uploaded_file)

    # 按鈕觸發 API 請求
    if st.button("送出"):
        if not user_input.strip():
            st.error("請輸入問題")
        else:
            # 儲存使用者的問題到 session_state
            st.session_state["messages"].append({"role": "user", "content": user_input})

            # 如果有上傳檔案的內容，將其加入到問題中
            if file_content:
                st.session_state["messages"].append({"role": "user", "content": f"這是上傳的檔案內容：{file_content}"})

            # 設定 API 請求的 headers 和 payload
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }

            # Payload 用於 API 請求
            data = {
                "model": "gpt-4o-mini",
                "messages": st.session_state["messages"],
                "max_tokens": 10000
            }

            try:
                # 發送 API 請求
                response = requests.post(api_url, headers=headers, json=data)

                # 檢查請求是否成功
                if response.status_code == 200:
                    response_json = response.json()
                    answer = response_json['choices'][0]['message']['content']
                
                    # 儲存 AI 回應到 session_state
                    st.session_state["messages"].append({"role": "assistant", "content": answer})

                else:
                    st.error(f"API 調用失敗：狀態碼 {response.status_code}，錯誤訊息：{response.text}")

            except Exception as e:
                st.error(f"發生錯誤：{e}")

    # 顯示歷史對話串，反轉順序讓新的對話顯示在上方
    if st.session_state["messages"]:
        st.write("### 歷史對話：")
        for msg in reversed(st.session_state["messages"]):  # 反向顯示對話串
            if msg["role"] == "user":
                st.markdown(format_user_message(msg['content']), unsafe_allow_html=True)
            elif msg["role"] == "assistant":
                st.markdown(format_ai_message(msg['content']), unsafe_allow_html=True)




# ---------------------------------------------------------------------------------------------------------------
# HomeWork4分頁
elif st.session_state.selected_tab == "HomeWork4":
    # 請使用你的 API 金鑰
    OPENAI_API_KEY = 'sk-UebxWPd44CJnjYPWt3e85ogp7DOqXSj2AuGEUarE9NjUckni'

   openai.api_key = OPENAI_API_KEY

    st.title("與助理聊天")

    # 在側邊欄中設定助理的參數
    with st.sidebar:
        st.header("助理設定")
        assistant_id = st.text_input("請輸入助理 ID", "asst_abc123")
        new_name = st.text_input("新的助理名稱", "Helper")
        new_instructions = st.text_area("助理指令", "You are an assistant, and you have access to files to answer user questions.")
        model = st.selectbox("選擇模型", ["gpt-4o", "gpt-3.5-turbo"])
        temperature = st.slider("溫度設定", min_value=0.0, max_value=2.0, value=1.0)

        # 設定 max_tokens 預設值為 150
        max_tokens = st.number_input("最大 token 數", min_value=1, max_value=4096, value=150)

        # 檔案上傳
        uploaded_file = st.file_uploader("請上傳檔案", type=['txt', 'pdf', 'jsonl', 'csv'])

        # 按鈕用來發送請求更新助理
        if st.button("更新助理"):
            try:
                # 上傳檔案
                file_id = None
                if uploaded_file is not None:
                    # 將上傳的檔案傳送到 OpenAI
                    file_response = client.files.create(
                        file=uploaded_file,
                        purpose="assistants"
                    )
                    st.success("檔案上傳成功！")
                    file_id = file_response['id']  # 獲取檔案 ID

                # 發送更新請求到 OpenAI API
                my_updated_assistant = client.beta.assistants.update(
                    assistant_id,
                    instructions=new_instructions,
                    name=new_name,
                    model=model,
                    tools=[{"type": "file_search", "file_id": file_id}] if file_id else []  # 確保有檔案 ID 才添加
                )

                # 顯示更新後的助理信息
                st.write("助理更新成功！")
                st.json(my_updated_assistant)

            except Exception as e:
                st.error(f"更新助理時發生錯誤: {e}")

    # 聊天功能
    st.subheader("ChatGPT 聊天機器人")

    if "messages" not in st.session_state:
        st.session_state.messages = []

    # 顯示歷史消息
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # 如果有上傳的檔案，讀取內容並提供給助理分析
    if uploaded_file is not None:
        if uploaded_file.type == "text/plain":
            file_content = uploaded_file.read().decode("utf-8")
        elif uploaded_file.type == "application/pdf":
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            file_content = ""
            for page in pdf_reader.pages:
                file_content += page.extract_text()
        elif uploaded_file.type == "application/json":
            import json
            file_content = json.load(uploaded_file)
            file_content = json.dumps(file_content)  # 將 JSON 轉換為字串
        elif uploaded_file.type == "text/csv":
            # 如果是 CSV 檔案，使用 pandas 讀取
            file_content = pd.read_csv(uploaded_file).to_string(index=False)

        # 準備用於分析的訊息
        if uploaded_file is not None:
            try:
                # 將上傳的檔案傳送到 OpenAI
                file_response = client.files.create(
                    file=uploaded_file,
                    purpose="assistants"  # 設定目的為助理
                )
                st.success("檔案上傳成功！")
                file_id = file_response.get('id')  # 獲取檔案 ID
                st.write(f"檔案 ID: {file_id}")
            except Exception as e:
                st.error(f"上傳檔案時發生錯誤: {e}")
        else:
            st.warning("請選擇一個檔案上傳！")

    # 用戶輸入問題
    if prompt := st.chat_input("請輸入您的問題:"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            response = client.chat.completions.create(
                model=model,
                messages=st.session_state.messages,
                temperature=temperature,
                max_tokens=max_tokens  # 使用側邊欄設定的最大 token 數
            )
            full_response = response.choices[0].message.content.strip()
            st.markdown(full_response)

        st.session_state.messages.append({"role": "assistant", "content": full_response})
